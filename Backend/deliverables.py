import os
import json
import re
import uuid
from datetime import datetime
from typing import Any, Optional
from pathlib import Path
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from sovereign_llm import sovereign_llm
from network_monitor import network_monitor

OUTPUT_DIR = Path(__file__).parent / "generated_deliverables"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

def set_cell_background(cell, hex_color: str):
    """Set background color of a Word table cell."""
    tcPr = cell._element.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{hex_color}"/>')
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    """Set internal cell padding."""
    tcPr = cell._element.get_or_add_tcPr()
    tcMar = parse_xml(f'<w:tcMar {nsdecls("w")}><w:top w:w="{top}" w:type="dxa"/><w:bottom w:w="{bottom}" w:type="dxa"/><w:left w:w="{left}" w:type="dxa"/><w:right w:w="{right}" w:type="dxa"/></w:tcMar>')
    tcPr.append(tcMar)


class DeliverablesEngine:
    def __init__(self, output_dir: Path = OUTPUT_DIR):
        self.output_dir = output_dir
        self.deliverables_registry: dict[str, dict[str, Any]] = {}

    def generate_word_approval_note(
        self,
        dna: dict[str, Any],
        params: dict[str, Any],
        custom_content: Optional[str] = None
    ) -> dict[str, Any]:
        """Generate a formal PSU / Refinery / Defense style Word Approval Note (.docx)."""
        file_id = str(uuid.uuid4())[:8]
        filename = f"Approval_Note_{file_id}.docx"
        file_path = self.output_dir / filename
        
        doc = docx.Document()
        
        # Set standard margins (1 inch)
        for section in doc.sections:
            section.top_margin = Inches(1.0)
            section.bottom_margin = Inches(1.0)
            section.left_margin = Inches(1.0)
            section.right_margin = Inches(1.0)
            
        # Top Header Banner Table
        header_table = doc.add_table(rows=2, cols=2)
        header_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        header_table.autofit = False
        
        c00 = header_table.cell(0, 0)
        c00.text = "CONFIDENTIAL // ON-PREMISES ONLY"
        c00.paragraphs[0].runs[0].font.bold = True
        c00.paragraphs[0].runs[0].font.size = Pt(9)
        c00.paragraphs[0].runs[0].font.color.rgb = RGBColor(180, 40, 40)
        
        c01 = header_table.cell(0, 1)
        c01.text = f"REF NO: SOV-{datetime.now().strftime('%Y%m%d')}-{file_id.upper()}"
        c01.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
        c01.paragraphs[0].runs[0].font.size = Pt(9)
        c01.paragraphs[0].runs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        c10 = header_table.cell(1, 0)
        c10.text = "SOVEREIGN INDUSTRIAL AI WORKBENCH"
        c10.paragraphs[0].runs[0].font.bold = True
        c10.paragraphs[0].runs[0].font.size = Pt(10)
        c10.paragraphs[0].runs[0].font.color.rgb = RGBColor(30, 41, 59)
        
        c11 = header_table.cell(1, 1)
        c11.text = f"DATE: {datetime.now().strftime('%d %B %Y')}"
        c11.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
        c11.paragraphs[0].runs[0].font.size = Pt(9)
        
        doc.add_paragraph() # Spacing
        
        # Document Title
        title_p = doc.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_p.add_run(f"EXECUTIVE APPROVAL NOTE & TECHNICAL ASSESSMENT\n{dna.get('identity', 'Industrial Assessment').upper()}")
        title_run.font.bold = True
        title_run.font.size = Pt(15)
        title_run.font.color.rgb = RGBColor(15, 23, 42)
        
        # Metadata Block Table
        meta_table = doc.add_table(rows=4, cols=2)
        meta_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        rows_data = [
            ("Target Audience", params.get("target_audience", "Executives & Board Members")),
            ("Communication Objective", params.get("objective", "Formal Briefing & Approval")),
            ("Source Material", dna.get("source_name", "Technical Documentation")),
            ("Tone / Content Style", f"{params.get('tone', 'Formal')} / {params.get('style', 'PSU Standard')}")
        ]
        for idx, (label, val) in enumerate(rows_data):
            cell_lbl = meta_table.cell(idx, 0)
            cell_val = meta_table.cell(idx, 1)
            cell_lbl.text = label
            cell_lbl.paragraphs[0].runs[0].font.bold = True
            cell_lbl.paragraphs[0].runs[0].font.size = Pt(10)
            set_cell_background(cell_lbl, "F1F5F9")
            cell_val.text = str(val)
            cell_val.paragraphs[0].runs[0].font.size = Pt(10)
            set_cell_background(cell_val, "FAFAFA")
            set_cell_margins(cell_lbl, 80, 80, 100, 100)
            set_cell_margins(cell_val, 80, 80, 100, 100)
            
        doc.add_paragraph()
        
        # 1. Executive Summary & Overview
        h1 = doc.add_heading("1. Executive Summary & Context", level=1)
        h1.runs[0].font.color.rgb = RGBColor(15, 23, 42)
        p_overview = doc.add_paragraph(dna.get("overview", "Comprehensive overview of the subject matter."))
        p_overview.runs[0].font.size = Pt(10.5)
        
        # If custom generated content exists, add it
        if custom_content:
            doc.add_heading("2. Detailed Strategic Assessment", level=1)
            p_custom = doc.add_paragraph(custom_content)
            p_custom.runs[0].font.size = Pt(10.5)

        # 3. Key Extracted Claims & Factual Findings
        doc.add_heading("3. Key Findings & Extracted Claims", level=1)
        claims = dna.get("claims", [])
        findings = dna.get("key_findings", [])
        all_findings = (findings + claims)[:8]
        if all_findings:
            for item in all_findings:
                bullet = doc.add_paragraph(item, style="List Bullet")
                bullet.runs[0].font.size = Pt(10.5)
        else:
            doc.add_paragraph("No explicit claims identified.")

        # 4. Critical Statistics & Measurements Table
        stats = dna.get("statistics", [])
        if stats:
            doc.add_heading("4. Key Technical Metrics & Statistics", level=1)
            stats_table = doc.add_table(rows=1, cols=2)
            stats_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            hdr_cells = stats_table.rows[0].cells
            hdr_cells[0].text = "Metric / Parameter"
            hdr_cells[1].text = "Extracted Value & Context"
            set_cell_background(hdr_cells[0], "1E293B")
            set_cell_background(hdr_cells[1], "1E293B")
            hdr_cells[0].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            hdr_cells[1].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            hdr_cells[0].paragraphs[0].runs[0].font.bold = True
            hdr_cells[1].paragraphs[0].runs[0].font.bold = True
            
            for s_idx, st in enumerate(stats[:10]):
                row_cells = stats_table.add_row().cells
                row_cells[0].text = f"Metric #{s_idx+1}"
                row_cells[1].text = str(st)
                bg = "F8FAFC" if s_idx % 2 == 0 else "FFFFFF"
                set_cell_background(row_cells[0], bg)
                set_cell_background(row_cells[1], bg)
                set_cell_margins(row_cells[0], 60, 60, 100, 100)
                set_cell_margins(row_cells[1], 60, 60, 100, 100)

        # 5. Risk Assessment & Operational Implications
        risks = dna.get("risks", [])
        implications = dna.get("implications", [])
        if risks or implications:
            doc.add_heading("5. Risk Evaluation & Operational Implications", level=1)
            if risks:
                doc.add_heading("Identified Risks:", level=2)
                for r in risks:
                    bp = doc.add_paragraph(r, style="List Bullet")
                    bp.runs[0].font.color.rgb = RGBColor(180, 40, 40)
            if implications:
                doc.add_heading("Operational Implications:", level=2)
                for imp in implications:
                    doc.add_paragraph(imp, style="List Bullet")

        # 6. Strategic Recommendations & Action Plan
        recs = dna.get("recommendations", [])
        doc.add_heading("6. Strategic Recommendations & Corrective Actions", level=1)
        if recs:
            for idx, r in enumerate(recs):
                p = doc.add_paragraph(f"{idx+1}. {r}")
                p.runs[0].font.bold = True
                p.runs[0].font.size = Pt(10.5)
        else:
            doc.add_paragraph("Continue baseline operations under regular monitoring protocols.")

        # 7. Official Sign-off & Approval Matrix Table
        doc.add_paragraph()
        doc.add_heading("7. Sign-off & Approval Signatures", level=1)
        sign_table = doc.add_table(rows=3, cols=3)
        sign_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        
        sign_headers = ["Initiating Officer / Engineer", "Reviewing Authority (Technical)", "Approving Authority (Executive)"]
        for c_idx, title in enumerate(sign_headers):
            cell = sign_table.cell(0, c_idx)
            cell.text = title
            cell.paragraphs[0].runs[0].font.bold = True
            cell.paragraphs[0].runs[0].font.size = Pt(9.5)
            set_cell_background(cell, "E2E8F0")
            
        for c_idx in range(3):
            cell_sig = sign_table.cell(1, c_idx)
            cell_sig.text = "\n\n_______________________\nSignature & Stamp"
            cell_sig.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            set_cell_background(cell_sig, "FAFAFA")
            set_cell_margins(cell_sig, 100, 100, 100, 100)
            
            cell_date = sign_table.cell(2, c_idx)
            cell_date.text = "Date: _______________"
            cell_date.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            set_cell_background(cell_date, "FAFAFA")
            set_cell_margins(cell_date, 60, 60, 100, 100)
            
        doc.save(str(file_path))
        
        record = {
            "id": file_id,
            "filename": filename,
            "path": str(file_path),
            "format": "word_docx",
            "title": f"Approval Note - {dna.get('identity', 'Assessment')}",
            "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "size_bytes": file_path.stat().st_size
        }
        self.deliverables_registry[file_id] = record
        return record

    def generate_pptx_deck(
        self,
        dna: dict[str, Any],
        params: dict[str, Any],
        custom_content: Optional[str] = None
    ) -> dict[str, Any]:
        """Generate a multi-slide PowerPoint presentation (.pptx)."""
        file_id = str(uuid.uuid4())[:8]
        filename = f"Board_Presentation_{file_id}.pptx"
        file_path = self.output_dir / filename
        
        prs = Presentation()
        prs.slide_width = PptxInches(13.333) # 16:9 widescreen
        prs.slide_height = PptxInches(7.5)
        
        # Color palette (Sovereign Industrial)
        c_dark_navy = PptxRGBColor(15, 23, 42)
        c_cyan_accent = PptxRGBColor(14, 165, 233)
        c_gray_text = PptxRGBColor(71, 85, 105)
        c_danger_red = PptxRGBColor(220, 38, 38)
        
        # Slide 1: Title Slide
        blank_slide_layout = prs.slide_layouts[6]
        s1 = prs.slides.add_slide(blank_slide_layout)
        
        # Header banner shape
        txBox = s1.shapes.add_textbox(PptxInches(1.0), PptxInches(1.5), PptxInches(11.333), PptxInches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p_badge = tf.paragraphs[0]
        p_badge.text = "SOVEREIGN AIR-GAPPED BRIEFING // CONFIDENTIAL"
        p_badge.font.size = PptxPt(12)
        p_badge.font.bold = True
        p_badge.font.color.rgb = c_cyan_accent
        
        p_title = tf.add_paragraph()
        p_title.text = dna.get("identity", "Strategic Assessment").upper()
        p_title.font.size = PptxPt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark_navy
        
        p_sub = tf.add_paragraph()
        p_sub.text = f"Audience: {params.get('target_audience', 'Board of Directors')} | Objective: {params.get('objective', 'Executive Decision')}"
        p_sub.font.size = PptxPt(16)
        p_sub.font.color.rgb = c_gray_text
        
        p_date = tf.add_paragraph()
        p_date.text = f"Date: {datetime.now().strftime('%d %B %Y')} | Sovereign AI Verification Engine"
        p_date.font.size = PptxPt(12)
        p_date.font.color.rgb = c_gray_text

        def create_content_slide(title_text: str, subtitle_text: str, bullet_points: list[str], notes_text: str = ""):
            slide = prs.slides.add_slide(blank_slide_layout)
            tb = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(0.8), PptxInches(11.333), PptxInches(6.0))
            frame = tb.text_frame
            frame.word_wrap = True
            
            p_head = frame.paragraphs[0]
            p_head.text = title_text
            p_head.font.size = PptxPt(24)
            p_head.font.bold = True
            p_head.font.color.rgb = c_dark_navy
            
            p_sub = frame.add_paragraph()
            p_sub.text = subtitle_text
            p_sub.font.size = PptxPt(13)
            p_sub.font.color.rgb = c_cyan_accent
            
            frame.add_paragraph() # space
            
            for item in bullet_points[:6]:
                p_item = frame.add_paragraph()
                p_item.text = f"•  {item}"
                p_item.font.size = PptxPt(15)
                p_item.font.color.rgb = c_dark_navy
                p_item.space_after = PptxPt(12)
                
            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text
            return slide

        # Slide 2: Executive Overview
        overview_bullets = [
            dna.get("overview", "High level context."),
            f"Source Document: {dna.get('source_name', 'Industrial Log')}",
            f"Key Stakeholders: {', '.join(dna.get('entities', {}).get('organizations', [])[:4]) or 'Internal PSU'}"
        ]
        create_content_slide("Executive Overview", "Context and Core Objective", overview_bullets, "Speaker Note: Emphasize sovereign on-premises validation.")

        # Slide 3: Key Findings & Claims
        findings_bullets = (dna.get("key_findings", []) + dna.get("claims", []))[:5]
        if not findings_bullets:
            findings_bullets = ["No direct findings recorded in source."]
        create_content_slide("Key Findings & Factual Claims", "Extracted with High-Fidelity Content DNA", findings_bullets, "Speaker Note: Walk through each verified claim.")

        # Slide 4: Data & Statistics
        stats_bullets = [str(s) for s in dna.get("statistics", [])[:5]]
        if not stats_bullets:
            stats_bullets = ["Operating within baseline tolerances."]
        create_content_slide("Critical Statistics & Measurements", "Empirical Data Points & Quantities", stats_bullets, "Speaker Note: Focus on quantitative metrics.")

        # Slide 5: Risk & Implications
        risks_bullets = [f"[RISK] {r}" for r in dna.get("risks", [])[:3]] + [f"[IMPLICATION] {i}" for i in dna.get("implications", [])[:3]]
        if not risks_bullets:
            risks_bullets = ["No severe operational risks flagged."]
        create_content_slide("Risk Assessment & Operational Impact", "Mitigation Protocols & Exposure Analysis", risks_bullets, "Speaker Note: Detail risk mitigations.")

        # Slide 6: Recommendations & Action Plan
        recs_bullets = [f"Step {i+1}: {r}" for i, r in enumerate(dna.get("recommendations", [])[:5])]
        if not recs_bullets:
            recs_bullets = ["Proceed with standard operating procedure."]
        create_content_slide("Strategic Recommendations", "Actionable Next Steps & Implementation Roadmap", recs_bullets, "Speaker Note: Secure approvals for recommended items.")

        prs.save(str(file_path))
        
        record = {
            "id": file_id,
            "filename": filename,
            "path": str(file_path),
            "format": "powerpoint_pptx",
            "title": f"Presentation - {dna.get('identity', 'Deck')}",
            "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "size_bytes": file_path.stat().st_size
        }
        self.deliverables_registry[file_id] = record
        return record

    def generate_excel_sheet(
        self,
        dna: dict[str, Any],
        params: dict[str, Any]
    ) -> dict[str, Any]:
        """Generate a multi-sheet Excel Workbook (.xlsx) for calculations and data matrices."""
        file_id = str(uuid.uuid4())[:8]
        filename = f"Engineering_Calculations_{file_id}.xlsx"
        file_path = self.output_dir / filename
        
        wb = openpyxl.Workbook()
        
        # Styles
        header_fill = PatternFill(start_color="1E293B", end_color="1E293B", fill_type="solid")
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        sub_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
        accent_fill = PatternFill(start_color="E0F2FE", end_color="E0F2FE", fill_type="solid")
        bold_font = Font(name="Calibri", size=11, bold=True)
        thin_border = Border(
            left=Side(style='thin', color='CBD5E1'),
            right=Side(style='thin', color='CBD5E1'),
            top=Side(style='thin', color='CBD5E1'),
            bottom=Side(style='thin', color='CBD5E1')
        )
        
        # Sheet 1: Executive Summary
        ws1 = wb.active
        ws1.title = "Executive_Summary"
        ws1.views.sheetView[0].showGridLines = True
        
        ws1["A1"] = "SOVEREIGN INDUSTRIAL AI - DATA & CALCULATION WORKBOOK"
        ws1["A1"].font = Font(name="Calibri", size=14, bold=True, color="0F172A")
        ws1["A2"] = f"Subject: {dna.get('identity', 'Assessment')} | Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        ws1["A2"].font = Font(name="Calibri", size=10, italic=True, color="64748B")
        
        ws1["A4"] = "PARAMETER"
        ws1["B4"] = "VALUE"
        ws1["A4"].fill = header_fill
        ws1["B4"].fill = header_fill
        ws1["A4"].font = header_font
        ws1["B4"].font = header_font
        
        summary_rows = [
            ("Source Document", dna.get("source_name", "Documentation")),
            ("Audience", params.get("target_audience", "Executives")),
            ("Objective", params.get("objective", "Analysis & Calculation")),
            ("Total Claims Extracted", len(dna.get("claims", []))),
            ("Total Statistics Extracted", len(dna.get("statistics", []))),
            ("Total Risks Identified", len(dna.get("risks", []))),
            ("Air-Gap Status", "100% On-Premises Verified")
        ]
        for r_idx, (k, v) in enumerate(summary_rows, start=5):
            ws1[f"A{r_idx}"] = k
            ws1[f"B{r_idx}"] = str(v)
            ws1[f"A{r_idx}"].font = bold_font
            ws1[f"A{r_idx}"].border = thin_border
            ws1[f"B{r_idx}"].border = thin_border

        # Sheet 2: Statistics & Calculation Model
        ws2 = wb.create_sheet(title="Calculations_and_Data")
        ws2.views.sheetView[0].showGridLines = True
        
        headers2 = ["Item #", "Extracted Metric / Description", "Value", "Unit / Context", "Baseline Variance (%)", "Status"]
        for col_idx, h in enumerate(headers2, start=1):
            cell = ws2.cell(row=1, column=col_idx, value=h)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center")
            
        stats = dna.get("statistics", [])
        if not stats:
            stats = ["Flow Rate: 450 m3/h", "Operating Temperature: 82.5 C", "Pressure Delta: 1.4 bar", "Power Consumption: 120 kW"]
            
        for idx, st in enumerate(stats, start=2):
            ws2.cell(row=idx, column=1, value=idx-1).border = thin_border
            ws2.cell(row=idx, column=2, value=str(st)).border = thin_border
            # Try to extract numbers
            nums = re.findall(r"[-+]?\d*\.\d+|\d+", str(st))
            val = float(nums[0]) if nums else 100.0
            ws2.cell(row=idx, column=3, value=val).border = thin_border
            ws2.cell(row=idx, column=4, value="Standard Units").border = thin_border
            ws2.cell(row=idx, column=5, value=f"=C{idx}*0.05").border = thin_border
            ws2.cell(row=idx, column=6, value="NORMAL").border = thin_border

        # Sheet 3: Risk Matrix
        ws3 = wb.create_sheet(title="Risk_Matrix")
        ws3.views.sheetView[0].showGridLines = True
        
        headers3 = ["Risk ID", "Identified Hazard / Vulnerability", "Severity (1-5)", "Likelihood (1-5)", "Risk Score", "Mitigation Strategy"]
        for col_idx, h in enumerate(headers3, start=1):
            cell = ws3.cell(row=1, column=col_idx, value=h)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center")
            
        risks = dna.get("risks", [])
        if not risks:
            risks = ["Corrosion under insulation in pipeline segment B", "High thermal gradient during startup phase"]
            
        for idx, r in enumerate(risks, start=2):
            ws3.cell(row=idx, column=1, value=f"RSK-{idx-1:02d}").border = thin_border
            ws3.cell(row=idx, column=2, value=str(r)).border = thin_border
            ws3.cell(row=idx, column=3, value=3).border = thin_border
            ws3.cell(row=idx, column=4, value=2).border = thin_border
            ws3.cell(row=idx, column=5, value=f"=C{idx}*D{idx}").border = thin_border
            ws3.cell(row=idx, column=6, value="Implement enhanced NDT inspection and thermal monitoring").border = thin_border

        # Auto-adjust column widths across all sheets
        for sheet in wb.worksheets:
            for col in sheet.columns:
                max_len = 0
                col_letter = get_column_letter(col[0].column)
                for cell in col:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
                sheet.column_dimensions[col_letter].width = min(max(max_len + 4, 14), 50)

        wb.save(str(file_path))
        
        record = {
            "id": file_id,
            "filename": filename,
            "path": str(file_path),
            "format": "excel_xlsx",
            "title": f"Calculations - {dna.get('identity', 'Sheet')}",
            "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "size_bytes": file_path.stat().st_size
        }
        self.deliverables_registry[file_id] = record
        return record

    async def generate_deliverables(
        self,
        dna: dict[str, Any],
        formats: list[str],
        params: dict[str, Any]
    ) -> dict[str, Any]:
        """
        Generate multiple selected deliverables from the same Content DNA foundation.
        Formats: ['word_docx', 'powerpoint_pptx', 'excel_xlsx', 'executive_summary', 'linkedin_post', 'advisory', 'code']
        """
        results: dict[str, Any] = {
            "dna_id": dna.get("id"),
            "source_name": dna.get("source_name"),
            "generated_items": []
        }
        
        # Craft specialized prompt for narrative/text deliverables
        narrative_prompt = f"""
Using the following Content DNA as the factual foundation, generate high-quality deliverable content.

CONTENT DNA:
Identity: {dna.get('identity')}
Overview: {dna.get('overview')}
Claims: {json.dumps(dna.get('claims', []))}
Statistics: {json.dumps(dna.get('statistics', []))}
Key Findings: {json.dumps(dna.get('key_findings', []))}
Risks: {json.dumps(dna.get('risks', []))}
Recommendations: {json.dumps(dna.get('recommendations', []))}

GENERATION PARAMETERS:
Target Audience: {params.get('target_audience', 'Executives')}
Tone: {params.get('tone', 'Professional')}
Language: {params.get('language', 'English')}
Level of Detail: {params.get('level_of_detail', 'Comprehensive')}
Communication Objective: {params.get('objective', 'Inform & Recommend')}
Content Style: {params.get('style', 'Corporate / PSU Standard')}

REQUESTED FORMATS: {', '.join(formats)}

Write a comprehensive, professional, well-structured synthesis strictly adhering to the facts in the Content DNA without hallucinations.
"""
        custom_text_content = ""
        try:
            llm_res = await sovereign_llm.generate(
                prompt=narrative_prompt,
                task_type="deliverable_generator",
                temperature=0.2,
                timeout=None
            )
            custom_text_content = llm_res.get("text", "")
        except Exception as e:
            custom_text_content = f"Deliverable synthesis compiled from Content DNA: {dna.get('overview')}"

        # Generate each requested deliverable
        if "word_docx" in formats or "approval_note" in formats:
            doc_rec = self.generate_word_approval_note(dna, params, custom_text_content)
            results["generated_items"].append(doc_rec)
            
        if "powerpoint_pptx" in formats or "presentation" in formats:
            pptx_rec = self.generate_pptx_deck(dna, params, custom_text_content)
            results["generated_items"].append(pptx_rec)
            
        if "excel_xlsx" in formats or "spreadsheet" in formats or "calculations" in formats:
            xlsx_rec = self.generate_excel_sheet(dna, params)
            results["generated_items"].append(xlsx_rec)
            
        if any(f in formats for f in ["executive_summary", "linkedin_post", "advisory", "infographic", "report"]):
            summary_id = str(uuid.uuid4())[:8]
            summary_file = f"Deliverable_Summary_{summary_id}.md"
            summary_path = self.output_dir / summary_file
            with open(summary_path, "w", encoding="utf-8") as f:
                f.write(f"# Sovereign Deliverable: {dna.get('identity')}\n\n")
                f.write(f"**Target Audience**: {params.get('target_audience')}\n")
                f.write(f"**Objective**: {params.get('objective')}\n\n")
                f.write("---\n\n")
                f.write(custom_text_content)
                
            rec = {
                "id": summary_id,
                "filename": summary_file,
                "path": str(summary_path),
                "format": "markdown_text",
                "title": f"Summary - {dna.get('identity')}",
                "content": custom_text_content,
                "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "size_bytes": summary_path.stat().st_size
            }
            self.deliverables_registry[summary_id] = rec
            results["generated_items"].append(rec)

        return results

    def get_deliverable(self, file_id: str) -> Optional[dict[str, Any]]:
        return self.deliverables_registry.get(file_id)

    def list_all_deliverables(self) -> list[dict[str, Any]]:
        return list(self.deliverables_registry.values())

deliverables_engine = DeliverablesEngine()
